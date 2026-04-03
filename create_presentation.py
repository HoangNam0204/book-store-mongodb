from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = Path(__file__).resolve().parent / "Bao_cao_do_an_Book_Store.pptx"

MAROON = RGBColor(133, 0, 0)
MAROON_DARK = RGBColor(90, 0, 0)
GOLD = RGBColor(218, 173, 86)
TEXT = RGBColor(44, 44, 44)
MUTED = RGBColor(98, 98, 98)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(250, 245, 240)


def add_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.55))
    header.fill.solid()
    header.fill.fore_color.rgb = MAROON
    header.line.color.rgb = MAROON

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.55), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD


def add_footer(slide, text="Book Store | Lập trình Web nâng cao"):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(8.8), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = MAROON_DARK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.8), Inches(10.2), Inches(0.5))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Arial"
        sp.font.size = Pt(15)
        sp.font.color.rgb = MUTED

    add_footer(slide)


def add_bullets_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = RGBColor(230, 220, 214)

    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.45), Inches(11.2), Inches(3.95))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    for index, item in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = "Arial"
        p.font.size = Pt(20 if level == 0 else 15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10 if level == 0 else 4)
        if level == 0:
            p.font.bold = True if text.endswith(":") else False

    return slide


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.0), Inches(11.9), Inches(5.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = LIGHT
    banner.line.color.rgb = RGBColor(230, 220, 214)

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(10.8), Inches(1.3))
    p1 = title_box.text_frame.paragraphs[0]
    p1.text = "BÁO CÁO ĐỒ ÁN MÔN HỌC"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = MAROON_DARK
    p1.alignment = PP_ALIGN.CENTER

    p2 = title_box.text_frame.add_paragraph()
    p2.text = "LẬP TRÌNH WEB NÂNG CAO"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = MAROON
    p2.alignment = PP_ALIGN.CENTER

    project_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(10.6), Inches(1.0))
    pp = project_box.text_frame.paragraphs[0]
    pp.text = "Đề tài: Xây dựng website bán sách trực tuyến Book Store"
    pp.font.size = Pt(24)
    pp.font.bold = True
    pp.font.color.rgb = TEXT
    pp.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(1.6), Inches(4.0), Inches(9.8), Inches(1.8))
    tf = info_box.text_frame
    lines = [
        "Sinh viên thực hiện: [Điền họ tên thành viên / nhóm]",
        "Lớp: [Điền lớp]",
        "GVHD: [Điền tên giảng viên]",
        "Trường: [Điền tên trường]",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

    add_footer(slide, "Book Store | Báo cáo đồ án môn học")


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullets_slide(
        prs,
        "1. Giới thiệu đề tài",
        [
            "Tên đề tài: Website bán sách trực tuyến Book Store",
            "Mục tiêu xây dựng một hệ thống web hoàn chỉnh phục vụ bán sách và quản lý dữ liệu trực tuyến.",
            "Các chức năng chính:",
            ("Xem danh sách sản phẩm, tìm kiếm và lọc theo thể loại hoặc giá", 1),
            ("Đăng ký, đăng nhập, giỏ hàng, thanh toán, lịch sử đơn hàng", 1),
            ("Trang quản trị cho admin để CRUD sản phẩm, danh mục và quản lý đơn hàng", 1),
            ("Dữ liệu được lưu trữ bằng MongoDB và triển khai online bằng Render", 1),
        ],
        "Tổng quan về mục tiêu và phạm vi của đồ án",
    )

    add_bullets_slide(
        prs,
        "2. HTML",
        [
            "HTML là ngôn ngữ đánh dấu dùng để xây dựng cấu trúc trang web.",
            "Trong đề tài, HTML được dùng để tạo bố cục cho các trang:",
            ("Trang chủ, danh sách sản phẩm, chi tiết sản phẩm", 1),
            ("Giỏ hàng, thanh toán, lịch sử đơn hàng", 1),
            ("Đăng nhập, đăng ký, liên hệ và quản trị", 1),
            "HTML giúp tổ chức nội dung rõ ràng, dễ kết hợp với CSS và JavaScript.",
        ],
        "Ngôn ngữ nền tảng tạo cấu trúc giao diện",
    )

    add_bullets_slide(
        prs,
        "3. CSS",
        [
            "CSS dùng để định dạng giao diện và tạo trải nghiệm trực quan cho người dùng.",
            "Vai trò trong đồ án:",
            ("Thiết kế menu, banner, card sản phẩm, form và footer", 1),
            ("Tạo màu sắc đồng bộ theo phong cách của website Book Store", 1),
            ("Tùy chỉnh bố cục responsive trên nhiều kích thước màn hình", 1),
            "CSS giúp website đẹp hơn, dễ sử dụng hơn và chuyên nghiệp hơn khi demo.",
        ],
        "Công nghệ định dạng và hoàn thiện giao diện",
    )

    add_bullets_slide(
        prs,
        "4. JavaScript",
        [
            "JavaScript là ngôn ngữ xử lý tương tác trên giao diện web.",
            "Trong đề tài, JavaScript đảm nhiệm:",
            ("Tìm kiếm, lọc sản phẩm và cập nhật danh sách hiển thị", 1),
            ("Thêm vào giỏ hàng, mua ngay, cập nhật số lượng", 1),
            ("Gửi dữ liệu đăng nhập, đăng ký, liên hệ và thanh toán", 1),
            ("Hiển thị thông báo sau khi thêm, xóa hoặc cập nhật dữ liệu", 1),
        ],
        "Ngôn ngữ tạo tính tương tác cho website",
    )

    add_bullets_slide(
        prs,
        "5. Node.js",
        [
            "Node.js là môi trường chạy JavaScript phía server.",
            "Vai trò trong đồ án:",
            ("Xây dựng backend cho website", 1),
            ("Xử lý request từ client và phản hồi dữ liệu", 1),
            ("Kết nối với MongoDB qua Mongoose", 1),
            ("Quản lý session, giỏ hàng, đơn hàng và xác thực người dùng", 1),
            "Nhờ Node.js, website hoạt động như một ứng dụng web hoàn chỉnh thay vì chỉ là các trang tĩnh.",
        ],
        "Môi trường backend chính của hệ thống",
    )

    add_bullets_slide(
        prs,
        "6. MongoDB",
        [
            "MongoDB là hệ quản trị cơ sở dữ liệu NoSQL lưu trữ dữ liệu dưới dạng document.",
            "Trong đồ án, MongoDB lưu các collection quan trọng:",
            ("users", 1),
            ("products", 1),
            ("categories", 1),
            ("carts", 1),
            ("orders", 1),
            ("contactmessages", 1),
            "MongoDB linh hoạt, phù hợp với dữ liệu web hiện đại và dễ tích hợp với Node.js.",
        ],
        "Cơ sở dữ liệu chính của hệ thống",
    )

    add_bullets_slide(
        prs,
        "7. MongoDB Atlas và MongoDB Compass",
        [
            "MongoDB Atlas được dùng để host cơ sở dữ liệu online, giúp website hoạt động 24/7 trên Internet.",
            "MongoDB Compass được dùng để quản lý và quan sát dữ liệu trực quan.",
            "Ý nghĩa khi demo:",
            ("Thao tác trên website online", 1),
            ("Refresh trong Compass", 1),
            ("Quan sát dữ liệu thay đổi theo thời gian thực", 1),
            "Điều này giúp chứng minh hệ thống web đã kết nối đúng với database online.",
        ],
        "Kết hợp giữa database online và công cụ quản lý trực quan",
    )

    add_bullets_slide(
        prs,
        "8. ExpressJS",
        [
            "ExpressJS là framework của Node.js dùng để xây dựng server và API.",
            "Vai trò trong đề tài:",
            ("Tạo route cho frontend và backend", 1),
            ("Xử lý API đăng nhập, sản phẩm, giỏ hàng, đơn hàng, admin", 1),
            ("Kết nối frontend với MongoDB qua tầng xử lý backend", 1),
            "ExpressJS giúp cấu trúc project rõ ràng và triển khai nhanh các chức năng CRUD.",
        ],
        "Framework backend hỗ trợ xây dựng route và API",
    )

    add_bullets_slide(
        prs,
        "9. ReactJS",
        [
            "ReactJS là thư viện JavaScript dùng để xây dựng giao diện theo component.",
            "Trong đồ án hiện tại, nhóm chưa triển khai frontend bằng ReactJS.",
            "Tuy nhiên ReactJS là hướng phát triển phù hợp trong tương lai:",
            ("Tách giao diện thành component dễ tái sử dụng", 1),
            ("Quản lý trạng thái hiệu quả hơn", 1),
            ("Dễ mở rộng khi hệ thống lớn hơn", 1),
            "Vì vậy ReactJS được xem là công nghệ tham khảo và định hướng nâng cấp cho dự án.",
        ],
        "Công nghệ tham khảo và hướng phát triển tiếp theo",
    )

    add_bullets_slide(
        prs,
        "10. Demo hệ thống",
        [
            "Các nội dung sẽ demo khi thuyết trình:",
            ("Mở website online trên Render", 1),
            ("Đăng nhập tài khoản customer và thêm sản phẩm vào giỏ hàng", 1),
            ("Thanh toán và tạo đơn hàng", 1),
            ("Mở MongoDB Compass và refresh collection orders hoặc carts", 1),
            ("Đăng nhập admin, thêm hoặc xóa sản phẩm và danh mục", 1),
            ("Refresh collection products hoặc categories để thấy dữ liệu thay đổi", 1),
            "Link demo online: https://book-store-mongodb.onrender.com/Bookstore/Trang_chu.html",
        ],
        "Kịch bản trình diễn chức năng chính của website",
    )

    add_bullets_slide(
        prs,
        "11. Tài liệu tham khảo",
        [
            "Node.js: https://nodejs.org",
            "ExpressJS: https://expressjs.com",
            "MongoDB Docs: https://www.mongodb.com/docs",
            "MongoDB Compass: https://www.mongodb.com/products/compass",
            "MongoDB Atlas: https://www.mongodb.com/atlas",
            "Render Docs: https://render.com/docs",
            "GitHub Docs: https://docs.github.com",
            "Giáo trình và tài liệu môn Lập trình Web nâng cao",
        ],
        "Nguồn tài liệu được sử dụng trong quá trình thực hiện đồ án",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.45), Inches(11.1), Inches(4.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = RGBColor(230, 220, 214)

    thanks = slide.shapes.add_textbox(Inches(1.4), Inches(2.3), Inches(10.5), Inches(1.4))
    p = thanks.text_frame.paragraphs[0]
    p.text = "Xin chân thành cảm ơn quý thầy cô\nvà các bạn đã lắng nghe"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = MAROON_DARK

    note = thanks.text_frame.add_paragraph()
    note.text = "Nhóm em rất mong nhận được ý kiến đóng góp để hoàn thiện đề tài hơn."
    note.alignment = PP_ALIGN.CENTER
    note.font.size = Pt(17)
    note.font.color.rgb = MUTED

    add_footer(slide, "Book Store | Xin cảm ơn")

    prs.save(OUTFILE)
    print(f"Created: {OUTFILE}")


if __name__ == "__main__":
    build_deck()
